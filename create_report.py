"""Create_report.ipynb

Automatically generated by Colaboratory.

Original file is located at
    https://colab.research.google.com/drive/1us-RACqWKZThlMczTMP4TnB_5rgoQTXh
"""

import json
import os
import streamlit as st
from zipfile import ZipFile
from io import BytesIO
import datetime
import codecs
import base64
from docx import Document
from PIL import Image
from docx.shared import Inches
from docx.shared import Pt
from pptx import Presentation
from PIL import Image


# Create title for streamlit web app
st.title('Auto Report')
st.write("Make life more easier")
image = Image.open('thuy_si.jpg')
st.image(image)


# Get name of image from autoMl notebook
list_chart = ['Score Distribution', 'ROC', 'Cumulative Gain Chart', 'Lift chart',
              'Cumulative lift chart', 'Precision-Recall curve']


# Choose use case
USE_CASE = st.sidebar.selectbox('Use Case',('3k3d','50Mb'))
SLIDE_KQ_NAME = {'3k3d':'Model du bao thue bao dat 3k3d',
                 '50MB':'Model du bao thue bao khong dat 50MB'}


def delete_paragraph(paragraph):
    p = paragraph._element
    p.getparent().remove(p)
    p._p = p._element = None

def choose_bm6a_base():
    '''
    Choose BM6A template
    '''
    bm6a_3k3d = Document('BM06a_3k3d.docx')
    bm6a_50mb = Document('BM06a_3k3d.docx')
    bm6a_dict = {'3k3d': bm6a_3k3d, '50Mb': bm6a_50mb}
    bm6a_base = bm6a_dict[USE_CASE]
    return bm6a_base

def get_style_name(document):
    for style in bm6a_base.styles:
      print("style.name == %s" % style.name)

def create_custom_style(document, style_name: str):
  '''
  Create custome style
  '''
  obj_styles = document.styles
  obj_charstyle = obj_styles.add_style(style_name, WD_STYLE_TYPE.CHARACTER)
  obj_font = obj_charstyle.font
  obj_font.size = Pt(12)
  obj_font.name = 'Times New Roman'

def get_image_shape(document):
    for s in document.inline_shapes:
        print (s.height.cm,s.width.cm)

# Insert image into BM6a
def create_bm6a(bm6a_base):
    '''
    Firstly, delete all paragraph between part 3 and 4.
    Then add image to it
    '''
    section_index = 0

    for para in bm6a_base.paragraphs:
        if para.style.name == 'Heading 1':
           section_index += 1
        if section_index == 3:
           delete_paragraph(para)
        if section_index == 4:
           for chart_index in range(len(list_chart)):
             para = para.insert_paragraph_before('')
             r = para.add_run().add_picture(f'{list_chart[chart_index]}.png',
                                            width = Inches(7) , height=Inches(3))
             para = para.insert_paragraph_before(f'{list_chart[chart_index]}',style= None)
           para.insert_paragraph_before('Kết quả mô hình', style='Heading 1')
           bm6a_base.save('BM6a.docx')
           break
#create_bm6a(bm6a_base)

def create_slide_kq():
    X = Presentation()
    Layout = X.slide_layouts[0]
    first_slide = X.slides.add_slide(Layout)
    first_slide.shapes.title.text = "SLIDE BÁO CÁO"
    first_slide.placeholders[1].text = SLIDE_KQ_NAME[USE_CASE]

    second_layout = X.slide_layouts[5]
    for chart_index in range(len(list_chart)):
        second_slide = X.slides.add_slide(second_layout)
        second_slide.shapes.title.text = list_chart[chart_index]
        second_slide.shapes.add_picture(f'{list_chart[chart_index]}.png',
                                        left = Inches(1),
                                        top = Inches(2),
                                        height = Inches(3),
                                        width = Inches(6))
    X.save('Slide_KQ.pptx')

def save_uploadedfile(uploadedfile):
    """
    Save streamlit upload file
    """
    with open("autoML_output.json","wb") as f:
        f.write(uploadedfile.getbuffer())
        
def download_all_file(list_file: list):
    '''
    Zip all file then create link for downloading
    '''
    current_month = datetime.datetime.today().strftime('%m_%Y')
    zip_obj = ZipFile(f'BM_6a_6b_slide_kq_{current_month}.zip', 'w')
    # Add multiple to zip
    for file in list_file:
        zip_obj.write(file)
    zip_obj.close()

    with open(f'BM_6a_6b_slide_kq_{current_month}.zip', 'rb') as f:
        bytes = f.read()
        b64 = base64.b64encode(bytes).decode()
        href = f"<a href=\"data:file/zip;base64,{b64}\" download='BM_6a_6b_slide_kq_{current_month}.zip'>\
        Download all files\
        </a>"
    st.sidebar.markdown(href, unsafe_allow_html=True)

# Upload and save output file of AutoML notebook
file_upload = st.sidebar.file_uploader("Upload AutoML file", type=["json"])

if file_upload is not None:
    save_uploadedfile(file_upload)
    # Get image from notebook output
    content = json.loads(codecs.open("autoML_output.json",'r','utf-8-sig').read())
    list_imgs = content['paragraphs'][4]['results']['msg']
    chart_index = 0
    for img in list_imgs:
        if img['type'] == 'IMG':
           imgData = base64.b64decode(img['data'])
           with open(f'{list_chart[chart_index]}.png','wb') as f:
               f.write(imgData)
           chart_index += 1
    # Create bm6a
    bm6a_base = choose_bm6a_base()
    create_bm6a(bm6a_base)
    # Create slide kq
    create_slide_kq()
    # Dowload all file
    download_all_file(['BM6a.docx', 'Slide_KQ.pptx'])

